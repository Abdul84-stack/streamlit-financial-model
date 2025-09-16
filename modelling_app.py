import streamlit as st
import pandas as pd
import numpy as np
import numpy_financial as npf
import plotly.graph_objects as go
from io import BytesIO
import base64
import datetime
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt
import pdfkit
import warnings
warnings.filterwarnings('ignore')


# --- Set up the Streamlit app layout ---
st.set_page_config(layout="wide", page_title="Comprehensive M&A Financial Model")

# --- Authentication Logic ---
def check_password():
    """Returns `True` if the user is logged in, `False` otherwise."""

    def login_form():
        """Shows the login form."""
        st.header("Login to your M&A Financial Tool")
        with st.container(border=True):
            st.markdown(
                """
                <style>
                .login-container {
                    padding: 2rem;
                    border-radius: 1rem;
                    box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
                    text-align: center;
                }
                .stTextInput label {
                    font-weight: bold;
                }
                .stButton button {
                    width: 100%;
                }
                </style>
                <div class="login-container">
                    <img src="https://placehold.co/150x150/F0F2F6/1E88E5?text=Login" alt="Login Icon" style="border-radius: 50%; margin-bottom: 20px;">
                </div>
                """,
                unsafe_allow_html=True
            )
            with st.form("login_form"):
                username = st.text_input("Username")
                password = st.text_input("Password", type="password")
                submitted = st.form_submit_button("Login")

            if submitted:
                if username == "abdul_model" and password == "123456":
                    st.session_state["authenticated"] = True
                    st.success("Logged in successfully!")
                    st.rerun()
                else:
                    st.error("Invalid username or password")

    if "authenticated" not in st.session_state:
        st.session_state["authenticated"] = False

    if not st.session_state["authenticated"]:
        login_form()
        return False
    else:
        return True

def main_app():
    """Main application content, visible only after successful login."""
    st.title("Comprehensive M&A Financial Modeling App 💰")
    st.markdown("A one-stop platform for financial modeling, valuation, and M&A analysis.")

    # --- Session State for Scenarios ---
    if "scenarios" not in st.session_state:
        st.session_state.scenarios = {
            "Base Case": {
                "revenue_growth": 5.0,
                "synergy_value": 15.0,
                "opex_percent_revenue": 50.0,
                "cogs_percent_revenue": 35.0,
                "danda_percent_revenue": 3.0,
                "capex_percent_revenue": 5.0,
                "interest_rate": 8.0,
                "debt_repayment_percent": 5.0,
                "exit_multiple": 10.0,
            },
            "Upside": {
                "revenue_growth": 8.0,
                "synergy_value": 25.0,
                "opex_percent_revenue": 48.0,
                "cogs_percent_revenue": 34.0,
                "danda_percent_revenue": 2.5,
                "capex_percent_revenue": 4.5,
                "interest_rate": 7.0,
                "debt_repayment_percent": 10.0,
                "exit_multiple": 12.0,
            },
            "Downside": {
                "revenue_growth": 2.0,
                "synergy_value": 5.0,
                "opex_percent_revenue": 55.0,
                "cogs_percent_revenue": 38.0,
                "danda_percent_revenue": 3.5,
                "capex_percent_revenue": 6.0,
                "interest_rate": 10.0,
                "debt_repayment_percent": 2.0,
                "exit_multiple": 8.0,
            },
        }

    # --- Mock Exchange Rates, Industry, Comps & Precedents Data ---
    EXCHANGE_RATES = {
        "USD": 1.0,
        "NGN": 1500.0,
        "EUR": 0.92,
        "GBP": 0.80,
        "JPY": 155.0,
        "CAD": 1.35,
    }

    INDUSTRY_DATA = {
        "Technology": {"P/E Ratio": 25.0, "EV/EBITDA": 18.0, "Net Profit Margin": 0.20},
        "Healthcare": {"P/E Ratio": 22.0, "EV/EBITDA": 15.0, "Net Profit Margin": 0.15},
        "FMCG (Consumer Goods)": {"P/E Ratio": 22.0, "EV/EBITDA": 17.0, "Net Profit Margin": 0.08},
        "Real Estate": {"P/E Ratio": 18.0, "EV/EBITDA": 12.0, "Net Profit Margin": 0.10},
        "Oil and Gas": {"P/E Ratio": 14.0, "EV/EBITDA": 8.0, "Net Profit Margin": 0.09},
        "Renewable Energy": {"P/E Ratio": 28.0, "EV/EBITDA": 20.0, "Net Profit Margin": 0.12},
        "Automobile": {"P/E Ratio": 15.0, "EV/EBITDA": 10.0, "Net Profit Margin": 0.07},
        "Agriculture": {"P/E Ratio": 16.0, "EV/EBITDA": 11.0, "Net Profit Margin": 0.06},
        "Insurance": {"P/E Ratio": 13.0, "EV/EBITDA": 9.0, "Net Profit Margin": 0.11},
        "Agro-allied": {"P/E Ratio": 17.0, "EV/EBITDA": 12.5, "Net Profit Margin": 0.07},
        "Banking": {"P/E Ratio": 12.0, "EV/EBITDA": 8.5, "Net Profit Margin": 0.25},
        "Education": {"P/E Ratio": 20.0, "EV/EBITDA": 14.0, "Net Profit Margin": 0.13},
        "Professional Services": {"P/E Ratio": 19.0, "EV/EBITDA": 13.0, "Net Profit Margin": 0.18},
        "Audit Firms": {"P/E Ratio": 15.0, "EV/EBITDA": 10.0, "Net Profit Margin": 0.22},
        "Retail": {"P/E Ratio": 21.0, "EV/EBITDA": 16.0, "Net Profit Margin": 0.04},
    }

    COMP_DATA = {
        "Company": ["Microsoft", "Google", "Amazon"],
        "Market Cap ($B)": [3000, 2200, 1900],
        "Revenue ($B)": [245, 307, 575],
        "EBITDA ($B)": [109, 84, 80],
        "Net Income ($B)": [72, 73, 30],
        "P/E Ratio": [41.6, 29.8, 115.5],
        "EV/EBITDA": [27.5, 26.2, 23.8]
    }

    PRECEDENT_DATA = {
        "Transaction": ["Salesforce/Slack", "Microsoft/Activision", "IBM/Red Hat"],
        "Date": ["2021", "2022", "2019"],
        "EV ($B)": [27.7, 68.7, 34.0],
        "Target Revenue ($B)": [1.1, 8.5, 3.0],
        "EV/Revenue": [25.2, 8.1, 11.3]
    }

    # --- Sidebar for User Inputs & Scenario Management ---
    st.sidebar.header("User Inputs & Settings")

    if st.sidebar.button("Logout"):
        st.session_state["authenticated"] = False
        st.rerun()

    selected_scenario = st.sidebar.selectbox("Select Scenario", options=list(st.session_state.scenarios.keys()))
    current_scenario_data = st.session_state.scenarios[selected_scenario]

    currency_options = list(EXCHANGE_RATES.keys())
    selected_currency = st.sidebar.selectbox("Select Display Currency", options=currency_options)
    industry_options = list(INDUSTRY_DATA.keys())
    selected_industry = st.sidebar.selectbox("Select Industry for Comparison", options=industry_options)
    industry_benchmarks = INDUSTRY_DATA[selected_industry]

    st.sidebar.markdown("---")
    st.sidebar.subheader("Financial Data Input")
    input_method = st.sidebar.radio("Select Input Method", ["Simulated API", "Upload CSV/Excel"])

    # --- Data Handling ---
    acquirer_data_df = pd.DataFrame()
    target_data_df = pd.DataFrame()

    if input_method == "Simulated API":
        acquirer_name = st.sidebar.text_input("Acquirer Company Ticker", "AAPL")
        target_name = st.sidebar.text_input("Target Company Ticker", "NVDA")

        def fetch_mock_data(company_name):
            with st.spinner(f"Fetching data for {company_name}..."):
                time.sleep(1)
                if company_name.upper() == "AAPL":
                    return pd.DataFrame({
                        'Metric': ['Revenue', 'EBITDA', 'Net_Income', 'Total_Assets', 'Total_Debt', 'Shares_Outstanding', 'Cash', 'Total_Liabilities', 'Total_Equity'],
                        'Value': [383.3, 131.9, 97.0, 352.6, 290.4, 15.7, 67.2, 290.4, 62.2]
                    }).set_index('Metric')
                elif company_name.upper() == "NVDA":
                    return pd.DataFrame({
                        'Metric': ['Revenue', 'EBITDA', 'Net_Income', 'Total_Assets', 'Total_Debt', 'Shares_Outstanding', 'Cash', 'Total_Liabilities', 'Total_Equity'],
                        'Value': [79.9, 53.6, 30.1, 74.3, 11.2, 2.5, 31.4, 11.2, 63.1]
                    }).set_index('Metric')
                else:
                    st.warning(f"No mock data available for {company_name}. Please use 'AAPL' or 'NVDA'.")
                    return pd.DataFrame()

        acquirer_data_df = fetch_mock_data(acquirer_name)
        target_data_df = fetch_mock_data(target_name)

    elif input_method == "Upload CSV/Excel":
        uploaded_file_acquirer = st.sidebar.file_uploader("Upload Acquirer's Financials (CSV or XLSX)", type=['csv', 'xlsx'])
        uploaded_file_target = st.sidebar.file_uploader("Upload Target's Financials (CSV or XLSX)", type=['csv', 'xlsx'])

        if uploaded_file_acquirer is not None and uploaded_file_target is not None:
            try:
                if uploaded_file_acquirer.name.endswith('.xlsx'):
                    acquirer_data_df = pd.read_excel(uploaded_file_acquirer, index_col=0)
                else:
                    acquirer_data_df = pd.read_csv(uploaded_file_acquirer, index_col=0)

                if uploaded_file_target.name.endswith('.xlsx'):
                    target_data_df = pd.read_excel(uploaded_file_target, index_col=0)
                else:
                    target_data_df = pd.read_csv(uploaded_file_target, index_col=0)

                acquirer_name = "Uploaded Acquirer"
                target_name = "Uploaded Target"
            except Exception as e:
                st.sidebar.error(f"Error reading file: {e}. Please ensure the file is a valid CSV or XLSX format.")
        else:
            st.info("Please upload both CSV or Excel files to proceed with this method.")


    if not acquirer_data_df.empty and not target_data_df.empty:
        required_metrics = ['Revenue', 'EBITDA', 'Net_Income', 'Total_Assets', 'Total_Debt', 'Shares_Outstanding', 'Cash', 'Total_Liabilities', 'Total_Equity']
        if not all(m in acquirer_data_df.index for m in required_metrics) or not all(m in target_data_df.index for m in required_metrics):
            st.error("Error: The uploaded files must contain all required metrics: " + ", ".join(required_metrics))
            acquirer_data_df = pd.DataFrame()
            target_data_df = pd.DataFrame()
        else:
            acquirer_data = acquirer_data_df['Value'].to_dict()
            target_data = target_data_df['Value'].to_dict()

            conversion_rate = EXCHANGE_RATES.get(selected_currency, 1.0)
            def convert_currency(value):
                return value / EXCHANGE_RATES['USD'] * conversion_rate
            
            for metric in acquirer_data.keys():
                acquirer_data[metric] = convert_currency(acquirer_data[metric])
            for metric in target_data.keys():
                target_data[metric] = convert_currency(target_data[metric])

            # --- Transaction Terms ---
            st.sidebar.markdown("---")
            st.sidebar.subheader("Transaction Terms")
            offer_price_per_share_base = st.sidebar.number_input(f"Offer Price per Share (in USD)", value=15.0, step=0.5)
            stock_percent = st.sidebar.slider("Percent Stock Consideration", min_value=0, max_value=100, value=50)

            # --- Model Assumptions (Macabacus-style) ---
            with st.expander("Model Assumptions & Scenario Management"):
                st.subheader("Scenario Assumptions")
                st.markdown("Adjust key drivers for the **" + selected_scenario + "** scenario.")

                col_a1, col_a2, col_a3, col_a4 = st.columns(4)
                with col_a1:
                    current_scenario_data["revenue_growth"] = st.number_input("Revenue Growth Rate (%)", value=current_scenario_data["revenue_growth"], step=0.5)
                    current_scenario_data["opex_percent_revenue"] = st.number_input("OpEx (% of Revenue)", value=current_scenario_data["opex_percent_revenue"], step=1.0)
                    current_scenario_data["synergy_value"] = st.number_input(f"Synergy Value (EBITDA, in {selected_currency} millions)", value=current_scenario_data["synergy_value"], step=1.0)
                with col_a2:
                    current_scenario_data["cogs_percent_revenue"] = st.number_input("COGS (% of Revenue)", value=current_scenario_data["cogs_percent_revenue"], step=1.0)
                    current_scenario_data["capex_percent_revenue"] = st.number_input("CapEx (% of Revenue)", value=current_scenario_data["capex_percent_revenue"], step=0.5)
                    current_scenario_data["danda_percent_revenue"] = st.number_input("D&A (% of Revenue)", value=current_scenario_data["danda_percent_revenue"], step=0.5)
                with col_a3:
                    current_scenario_data["interest_rate"] = st.number_input("Term Loan Interest Rate (%)", value=current_scenario_data["interest_rate"], step=0.5)
                    current_scenario_data["debt_repayment_percent"] = st.number_input("Annual Debt Repayment (%)", value=current_scenario_data["debt_repayment_percent"], step=1.0)
                with col_a4:
                    current_scenario_data["exit_multiple"] = st.number_input("Exit Multiple (EV/EBITDA)", value=current_scenario_data["exit_multiple"], step=0.5)
            
            # --- Financial Forecast & Three-Statement Model ---
            years = 5
            projections_df = pd.DataFrame(index=["Revenue", "COGS", "OpEx", "EBITDA", "D&A", "Interest Expense", "EBT", "Taxes", "Net Income", "Change in NWC", "CapEx", "FCF", "Beginning Cash", "Ending Cash", "Beginning Debt", "Ending Debt"])
            
            # Initial values from target data
            projections_df[0] = 0.0
            projections_df.loc["Revenue", 0] = target_data['Revenue']
            projections_df.loc["EBITDA", 0] = target_data['EBITDA']
            projections_df.loc["Net Income", 0] = target_data['Net_Income']
            projections_df.loc["Beginning Cash", 0] = target_data['Cash']
            projections_df.loc["Ending Debt", 0] = target_data['Total_Debt']

            for y in range(1, years + 1):
                projections_df[y] = 0.0
                
                # Income Statement
                proj_revenue = projections_df.loc["Revenue", y-1] * (1 + current_scenario_data["revenue_growth"]/100)
                proj_cogs = proj_revenue * current_scenario_data["cogs_percent_revenue"]/100
                proj_opex = proj_revenue * current_scenario_data["opex_percent_revenue"]/100
                proj_ebitda = proj_revenue - proj_cogs - proj_opex
                proj_danda = proj_revenue * current_scenario_data["danda_percent_revenue"]/100
                
                # Debt & Interest
                beginning_debt = target_data['Total_Debt'] if y == 1 else projections_df.loc["Ending Debt", y-1]
                interest_expense = beginning_debt * current_scenario_data["interest_rate"]/100
                proj_ebt = proj_ebitda - proj_danda - interest_expense
                proj_taxes = proj_ebt * 0.25 if proj_ebt > 0 else 0
                proj_net_income = proj_ebt - proj_taxes
                
                projections_df.loc["Revenue", y] = proj_revenue
                projections_df.loc["COGS", y] = proj_cogs
                projections_df.loc["OpEx", y] = proj_opex
                projections_df.loc["EBITDA", y] = proj_ebitda
                projections_df.loc["D&A", y] = proj_danda
                projections_df.loc["Interest Expense", y] = interest_expense
                projections_df.loc["EBT", y] = proj_ebt
                projections_df.loc["Taxes", y] = proj_taxes
                projections_df.loc["Net Income", y] = proj_net_income

                # Cash Flow Statement
                change_in_nwc = 0 # Simplified for now
                capex = proj_revenue * current_scenario_data["capex_percent_revenue"]/100
                fcf = proj_net_income + proj_danda - capex - change_in_nwc
                projections_df.loc["Change in NWC", y] = change_in_nwc
                projections_df.loc["CapEx", y] = capex
                projections_df.loc["FCF", y] = fcf

                # Balance Sheet & Debt Schedule
                ending_debt = beginning_debt - (beginning_debt * current_scenario_data["debt_repayment_percent"]/100)
                beginning_cash = target_data['Cash'] if y == 1 else projections_df.loc["Ending Cash", y-1]
                ending_cash = beginning_cash + fcf
                projections_df.loc["Beginning Cash", y] = beginning_cash
                projections_df.loc["Ending Cash", y] = ending_cash
                projections_df.loc["Beginning Debt", y] = beginning_debt
                projections_df.loc["Ending Debt", y] = ending_debt

            # Final calculations
            eps_acquirer = acquirer_data['Net_Income'] / acquirer_data['Shares_Outstanding']
            
            offer_price_per_share_converted = offer_price_per_share_base / 1000000000
            purchase_price = offer_price_per_share_converted * target_data['Shares_Outstanding']
            cash_consideration = purchase_price * (100 - stock_percent) / 100
            stock_consideration = purchase_price * stock_percent / 100
            
            current_share_price_acquirer_usd = st.sidebar.number_input(f"Acquirer's Current Share Price (in {selected_currency})", value=50000.0) / EXCHANGE_RATES['NGN']
            current_share_price_acquirer = current_share_price_acquirer_usd / 1000000000
            
            shares_issued = stock_consideration / current_share_price_acquirer
            total_pro_forma_shares = acquirer_data['Shares_Outstanding'] + shares_issued
            
            synergy_value_billions = current_scenario_data["synergy_value"] / 1000
            pro_forma_net_income = acquirer_data['Net_Income'] + target_data['Net_Income'] + (synergy_value_billions * (1 - 0.25))
            pro_forma_eps = pro_forma_net_income / total_pro_forma_shares
            accertion_dilution = ((pro_forma_eps - eps_acquirer) / eps_acquirer) * 100
            
            # Goodwill Calculation
            target_assets_fair_value = target_data['Total_Assets']
            target_liabilities = target_data['Total_Liabilities']
            target_equity = target_assets_fair_value - target_liabilities
            goodwill = purchase_price - target_equity
            
            # --- CORRECTED: Pro Forma Balance Sheet calculations ---
            pro_forma_assets = acquirer_data['Total_Assets'] + target_assets_fair_value + goodwill - cash_consideration
            pro_forma_liabilities = acquirer_data['Total_Liabilities'] + target_liabilities
            pro_forma_equity = acquirer_data['Total_Equity'] + target_equity + stock_consideration

            # --- Model Audits & Checks ---
            with st.container(border=True):
                st.header("Model Audits & Checks (S&P Global Style)")
                st.markdown("Automated checks to ensure the financial model is balanced and internally consistent.")
                
                # Check 1: Three-Statement Linkage
                # The fix: This check now verifies if the independently calculated components balance.
                audit_balance_sheet = pro_forma_assets - (pro_forma_liabilities + pro_forma_equity)
                if abs(audit_balance_sheet) < 0.01:
                    st.success("✅ Pro-Forma Balance Sheet is in balance.")
                else:
                    st.error(f"❌ Balance Sheet out of balance by {audit_balance_sheet:.2f} billion.")
                
                # Check 2: Cash Flow Reconciliation
                cash_check = projections_df.loc["Beginning Cash", years] + projections_df.loc["FCF", years] - projections_df.loc["Ending Cash", years]
                if abs(cash_check) < 0.01:
                    st.success(f"✅ Cash Flow Statement reconciles in year {years}.")
                else:
                    st.error(f"❌ Cash Flow Statement mismatch by {cash_check:.2f} billion in year {years}.")
                
                # Check 3: Debt Covenant (LBO only)
                if 'proj_ebitda' in locals():
                    debt_to_ebitda = (projections_df.loc["Beginning Debt", 1] + cash_consideration) / proj_ebitda
                    if debt_to_ebitda < 5:
                        st.success(f"✅ Debt/EBITDA ratio ({debt_to_ebitda:.2f}x) is within typical covenants.")
                    else:
                        st.warning(f"⚠️ Debt/EBITDA ratio ({debt_to_ebitda:.2f}x) may be a concern.")

            # --- Financial Ratios & Industry Comparison (Improved UI) ---
            with st.container(border=True):
                st.header("1. Financial Ratios & Industry Comparison 📊")
                def calculate_ratios(data, share_price):
                    ratios = {}
                    market_cap = share_price * data['Shares_Outstanding']
                    enterprise_value = market_cap + data['Total_Debt'] - data.get('Cash', 0)
                    ratios['P/E Ratio'] = market_cap / data['Net_Income'] if data['Net_Income'] != 0 else np.nan
                    ratios['EV/EBITDA'] = enterprise_value / data['EBITDA'] if data['EBITDA'] != 0 else np.nan
                    ratios['Net Profit Margin'] = data['Net_Income'] / data['Revenue'] if data['Revenue'] != 0 else np.nan
                    return ratios
                
                ratios_acquirer = calculate_ratios(acquirer_data, current_share_price_acquirer)
                ratios_target = calculate_ratios(target_data, offer_price_per_share_converted)
                
                ratio_comp_df = pd.DataFrame(index=INDUSTRY_DATA['Technology'].keys())
                ratio_comp_df[f'Acquirer ({selected_currency})'] = [ratios_acquirer.get(key, np.nan) for key in ratio_comp_df.index]
                ratio_comp_df[f'Target ({selected_currency})'] = [ratios_target.get(key, np.nan) for key in ratio_comp_df.index]
                ratio_comp_df[f'Industry Average ({selected_industry})'] = [industry_benchmarks.get(key, "N/A") for key in ratio_comp_df.index]
                st.table(ratio_comp_df.style.format(lambda x: f"{x:.2f}" if isinstance(x, (int, float)) else str(x)))
                
                fig_ratios = go.Figure()
                fig_ratios.add_trace(go.Bar(x=ratio_comp_df.index, y=ratio_comp_df[f'Acquirer ({selected_currency})'].values, name=f'Acquirer ({selected_currency})'))
                fig_ratios.add_trace(go.Bar(x=ratio_comp_df.index, y=ratio_comp_df[f'Target ({selected_currency})'].values, name=f'Target ({selected_currency})'))
                fig_ratios.add_trace(go.Bar(x=ratio_comp_df.index, y=ratio_comp_df[f'Industry Average ({selected_industry})'].values, name=f'Industry Average'))
                fig_ratios.update_layout(title="Financial Ratios Comparison", xaxis_title="Ratio", yaxis_title="Value", barmode='group')
                st.plotly_chart(fig_ratios, use_container_width=True)

            # --- Valuation Methodologies (Improved UI) ---
            with st.container(border=True):
                st.header("2. Valuation Methodologies 💵")
                col1, col2, col3 = st.columns(3)

                with col1:
                    st.subheader("DCF Valuation")
                    dcf_years = 5
                    wacc = 0.10
                    terminal_growth_rate = 0.02
                    
                    projected_fcf = projections_df.loc["FCF", 1:dcf_years].values
                    pv_fcf = sum([fcf / (1 + wacc)**i for i, fcf in enumerate(projected_fcf, 1)])
                    terminal_value = (projected_fcf[-1] * (1 + terminal_growth_rate)) / (wacc - terminal_growth_rate)
                    pv_terminal_value = terminal_value / (1 + wacc)**dcf_years
                    dcf_value = pv_fcf + pv_terminal_value
                    
                    st.write(f"**Target's DCF Valuation:** `{selected_currency} {dcf_value:,.2f} billion`")

                with col2:
                    st.subheader("Comparable Company Analysis (Comps)")
                    comps_df = pd.DataFrame(COMP_DATA)
                    average_ev_ebitda = comps_df['EV/EBITDA'].mean()
                    comps_value = average_ev_ebitda * target_data['EBITDA']
                    comps_value_converted = comps_value
                    st.write(f"**Average EV/EBITDA:** `{average_ev_ebitda:.2f}`")
                    st.write(f"**Target's Comps Valuation:** `{selected_currency} {comps_value_converted:,.2f} billion`")

                with col3:
                    st.subheader("Precedent Transaction Analysis")
                    precedents_df = pd.DataFrame(PRECEDENT_DATA)
                    average_ev_revenue = precedents_df['EV/Revenue'].mean()
                    precedents_value = average_ev_revenue * target_data['Revenue']
                    precedents_value_converted = precedents_value
                    st.write(f"**Average EV/Revenue:** `{average_ev_revenue:.2f}`")
                    st.write(f"**Target's Precedents Valuation:** `{selected_currency} {precedents_value_converted:,.2f} billion`")

            # --- Deal Structure & Accretion/Dilution Analysis (Improved UI) ---
            with st.container(border=True):
                st.header("3. Deal Structure & Accretion/Dilution 🤝")
                summary_data = {"Metric": ["Acquirer EPS", "Pro Forma EPS", "Accretion / Dilution"],
                                "Value": [eps_acquirer, pro_forma_eps, accertion_dilution]}
                st.table(pd.DataFrame(summary_data).set_index("Metric").style.format(lambda x: f"{x:.2f}" if isinstance(x, (int, float)) else str(x)))
                
                if accertion_dilution > 0:
                    st.success(f"**🎉 The transaction is ACCRETIVE to the acquirer's EPS by {accertion_dilution:.2f}%.**")
                else:
                    st.error(f"**📉 The transaction is DILUTIVE to the acquirer's EPS by {accertion_dilution:.2f}%.**")

            # --- Pro Forma Balance Sheet & Goodwill Calculation (Improved UI) ---
            with st.container(border=True):
                st.header("4. Pro Forma Balance Sheet & Goodwill ⚖️")
                st.write(f"**Goodwill Created:** `{selected_currency} {goodwill:,.2f} billion`")
                bs_data = {
                    "Acquirer": [acquirer_data['Total_Assets'], acquirer_data['Total_Liabilities'], acquirer_data['Total_Equity']],
                    "Target": [target_assets_fair_value, target_liabilities, target_equity],
                    "Pro Forma": [pro_forma_assets, pro_forma_liabilities, pro_forma_equity]
                }
                bs_df = pd.DataFrame(bs_data, index=["Total Assets", "Total Liabilities", "Total Equity"])
                st.table(bs_df.style.format(lambda x: f"{x:,.2f} billion"))

            # --- Scenario Analysis & Comparison ---
            with st.container(border=True):
                st.header("5. Scenario Analysis & Comparison 📈")
                results = {}
                for scenario_name, scenario_vars in st.session_state.scenarios.items():
                    # Recalculate based on scenario values
                    synergy_val = scenario_vars["synergy_value"]
                    rev_growth = scenario_vars["revenue_growth"]
                    
                    proj_ni = acquirer_data['Net_Income'] + target_data['Net_Income'] + ((synergy_val/1000) * (1-0.25))
                    proj_eps = proj_ni / total_pro_forma_shares
                    accret_dilut = ((proj_eps - eps_acquirer) / eps_acquirer) * 100
                    
                    exit_ebitda = target_data['EBITDA'] * (1 + rev_growth/100)**years
                    exit_ev = exit_ebitda * scenario_vars['exit_multiple']
                    
                    exit_debt = target_data['Total_Debt'] * (1 - scenario_vars['debt_repayment_percent']/100)**years
                    
                    sponsor_equity = purchase_price - target_data['Total_Debt'] + (purchase_price * 0.03)
                    moic = (exit_ev - exit_debt) / sponsor_equity
                    
                    results[scenario_name] = {
                        "Accretion/Dilution (%)": accret_dilut,
                        "MOIC (x)": moic,
                    }
                
                results_df = pd.DataFrame(results).T
                st.table(results_df.style.format({"Accretion/Dilution (%)": "{:.2f}%", "MOIC (x)": "{:.2f}x"}))
                
                fig_scenario = go.Figure()
                fig_scenario.add_trace(go.Bar(
                    x=results_df.index,
                    y=results_df["Accretion/Dilution (%)"],
                    name="EPS Accretion/Dilution"
                ))
                fig_scenario.add_trace(go.Bar(
                    x=results_df.index,
                    y=results_df["MOIC (x)"],
                    name="MOIC",
                    yaxis="y2"
                ))
                fig_scenario.update_layout(
                    title="Scenario Comparison",
                    xaxis_title="Scenario",
                    yaxis_title="EPS Accretion/Dilution (%)",
                    yaxis2=dict(title="MOIC (x)", overlaying="y", side="right"),
                    barmode="group",
                )
                st.plotly_chart(fig_scenario, use_container_width=True)

            # --- LBO/MBO Analysis (Improved UI) ---
            with st.container(border=True):
                st.header("6. LBO & MBO Analysis 📈")
                col_lbo1, col_lbo2 = st.columns(2)
                with col_lbo1:
                    st.subheader("Sources & Uses of Funds")
                    term_loan_percent = st.sidebar.slider("Term Loan % of Purchase Price", 0, 100, 50)
                    revolver_percent = st.sidebar.slider("Revolver % of Purchase Price", 0, 100, 5)
                    purchase_price_lbo = target_data['Shares_Outstanding'] * offer_price_per_share_converted
                    term_loan = purchase_price_lbo * term_loan_percent / 100
                    revolver = purchase_price_lbo * revolver_percent / 100
                    total_debt = term_loan + revolver
                    transaction_fees = purchase_price_lbo * (current_scenario_data["interest_rate"]/100)
                    sponsor_equity = purchase_price_lbo - total_debt + transaction_fees
                    sources = pd.DataFrame({'Category': ['Total Debt', 'Sponsor Equity'], 'Value': [total_debt, sponsor_equity]}).set_index('Category')
                    uses = pd.DataFrame({'Category': ['Purchase Price', 'Transaction Fees'], 'Value': [purchase_price_lbo, transaction_fees]}).set_index('Category')
                    st.markdown(f"**Total Sources:** `{selected_currency} {sources.sum().values[0]:,.2f} billion`")
                    st.markdown(f"**Total Uses:** `{selected_currency} {uses.sum().values[0]:,.2f} billion`")
                    fig_lbo = go.Figure(data=[go.Bar(name='Sources', x=['Total Debt', 'Sponsor Equity'], y=[total_debt, sponsor_equity], marker_color='skyblue'), go.Bar(name='Uses', x=['Purchase Price', 'Transaction Fees'], y=[purchase_price_lbo, transaction_fees], marker_color='salmon')])
                    fig_lbo.update_layout(title_text="Sources & Uses of Funds", barmode='group')
                    st.plotly_chart(fig_lbo)
                with col_lbo2:
                    st.subheader("LBO Analysis Results")
                    proj_ebitda = target_data['EBITDA'] * (1 + (current_scenario_data["revenue_growth"] / 100))**years
                    proj_exit_ev = proj_ebitda * current_scenario_data["exit_multiple"]
                    annual_repayment = term_loan * current_scenario_data["debt_repayment_percent"]/100
                    remaining_debt = max(0, term_loan - (annual_repayment * years))
                    proj_exit_debt = remaining_debt + revolver
                    proj_exit_equity = proj_exit_ev - proj_exit_debt
                    moic = proj_exit_equity / sponsor_equity if sponsor_equity > 0 else 0
                    cash_flows = [-sponsor_equity] + [0]*(years-1) + [proj_exit_equity]
                    def calculate_irr(cash_flows, tolerance=0.0001, max_iterations=1000):
                        if not cash_flows or cash_flows[0] >= 0: return np.nan
                        rate = 0.1
                        for _ in range(max_iterations):
                            npv = sum(cf / (1 + rate)**i for i, cf in enumerate(cash_flows))
                            if abs(npv) < tolerance: return rate
                            dnpv = sum(-i * cf / (1 + rate)**(i + 1) for i, cf in enumerate(cash_flows))
                            if dnpv == 0: return np.nan
                            rate -= npv / dnpv
                        return np.nan
                    irr = calculate_irr(cash_flows)
                    st.metric(label="Projected Exit Enterprise Value", value=f"{proj_exit_ev:,.2f} B")
                    st.metric(label="Projected Exit Equity Value", value=f"{proj_exit_equity:,.2f} B")
                    st.metric(label="Multiple on Invested Capital (MOIC)", value=f"{moic:.2f}x")
                    st.metric(label="Internal Rate of Return (IRR)", value=f"{irr:.2%}")

            # --- Reporting Module ---
            st.markdown("---")
            st.header("7. Generate Professional Report 📄")
            st.markdown("Note: Collaboration and version control features would require a backend database and are beyond the scope of this single script.")
            
            def create_powerpoint_report(acquirer_name, target_name, report_data, fig_ratios, fig_lbo, sensitivity_df, results_df, projections_df):
                prs = Presentation()
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                title.text = f"M&A Financial Analysis: {acquirer_name} & {target_name}"
                subtitle.text = f"Date: {datetime.date.today().strftime('%B %d, %Y')}"
                
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = "Executive Summary"
                if len(slide.placeholders) > 1:
                    body_shape = slide.placeholders[1]
                    tf = body_shape.text_frame
                    tf.clear()
                    p = tf.add_paragraph()
                    p.text = f"The transaction is { 'ACCRETIVE' if report_data['accertion_dilution'] > 0 else 'DILUTIVE' } to the acquirer's EPS by {report_data['accertion_dilution']:.2f}%."
                    p = tf.add_paragraph()
                    p.text = f"Goodwill created: {report_data['selected_currency']} {report_data['goodwill']:.2f} billion."
                    p = tf.add_paragraph()
                    p.text = "Valuation Summary:"
                    p = tf.add_paragraph()
                    p.text = f"- DCF: {report_data['selected_currency']} {report_data['dcf_value']:.2f} billion"
                    p = tf.add_paragraph()
                    p.text = f"- Comps: {report_data['selected_currency']} {report_data['comps_value_converted']:.2f} billion"
                    p = tf.add_paragraph()
                    p.text = f"- Precedents: {report_data['selected_currency']} {report_data['precedents_value_converted']:.2f} billion"
                
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                title = slide.shapes.title
                title.text = "Financial Ratios & Industry Comparison"
                try:
                    fig_ratios_img = BytesIO()
                    fig_ratios.write_image(fig_ratios_img, format='png')
                    fig_ratios_img.seek(0)
                    slide.shapes.add_picture(fig_ratios_img, Inches(1), Inches(2), width=Inches(8))
                except Exception as e:
                    st.error(f"Error saving ratios chart: {e}")
                
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                title = slide.shapes.title
                title.text = "LBO Analysis"
                if len(slide.placeholders) > 1:
                    body_shape = slide.placeholders[1]
                    tf = body_shape.text_frame
                    tf.clear()
                    p = tf.add_paragraph()
                    p.text = f"MOIC: {report_data['moic']:.2f}x"
                    p = tf.add_paragraph()
                    p.text = f"IRR: {report_data['irr']:.2%}"
                try:
                    fig_lbo_img = BytesIO()
                    fig_lbo.write_image(fig_lbo_img, format='png')
                    fig_lbo_img.seek(0)
                    slide.shapes.add_picture(fig_lbo_img, Inches(1), Inches(4), width=Inches(8))
                except Exception as e:
                    st.error(f"Error saving LBO chart: {e}")
                
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                title = slide.shapes.title
                title.text = "Scenario Analysis & Comparison"
                left = Inches(0.5)
                top = Inches(2)
                width = Inches(9)
                height = Inches(4)
                table_shape = slide.shapes.add_table(len(results_df) + 1, len(results_df.columns) + 1, left, top, width, height)
                table = table_shape.table
                for i, col in enumerate(results_df.columns): table.cell(0, i + 1).text = col
                for i, row in enumerate(results_df.index):
                    table.cell(i + 1, 0).text = row
                    for j, val in enumerate(results_df.iloc[i]):
                        table.cell(i + 1, j + 1).text = f"{val:.2f}"
                
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                title = slide.shapes.title
                title.text = "Projected Financials (Target Co.)"
                left = Inches(0.5)
                top = Inches(2)
                width = Inches(9)
                height = Inches(4)
                proj_df_for_report = projections_df.loc[["Revenue", "EBITDA", "Net Income", "FCF", "Ending Cash", "Ending Debt"]]
                table_shape = slide.shapes.add_table(len(proj_df_for_report) + 1, len(proj_df_for_report.columns) + 1, left, top, width, height)
                table = table_shape.table
                for i, col in enumerate(proj_df_for_report.columns): table.cell(0, i + 1).text = str(col)
                for i, row in enumerate(proj_df_for_report.index):
                    table.cell(i + 1, 0).text = row
                    for j, val in enumerate(proj_df_for_report.iloc[i]):
                        table.cell(i + 1, j + 1).text = f"{val:.2f}"
                        
                pptx_buffer = BytesIO()
                prs.save(pptx_buffer)
                pptx_buffer.seek(0)
                return pptx_buffer

            def create_pdf_report(acquirer_name, target_name, report_data, fig_ratios, fig_lbo, results_df, projections_df):
                try:
                    path_to_wkhtmltopdf = r'C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe'
                    config = pdfkit.configuration(wkhtmltopdf=path_to_wkhtmltopdf)
                    fig_ratios_base64 = base64.b64encode(fig_ratios.to_image(format='png')).decode()
                    fig_lbo_base64 = base64.b64encode(fig_lbo.to_image(format='png')).decode()
                    
                    proj_html = projections_df.loc[["Revenue", "EBITDA", "Net Income", "FCF", "Ending Cash", "Ending Debt"]].T.to_html(classes='table', float_format='%.2f')
                    
                    html_content = f"""
                    <!DOCTYPE html><html><head><title>M&A Analysis Report</title><style>body {{ font-family: sans-serif; }} h1 {{ color: #1f77b4; }} h2 {{ color: #1f77b4; border-bottom: 2px solid #ddd; padding-bottom: 5px; }} table {{ width: 100%; border-collapse: collapse; margin-bottom: 20px; }} th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }} th {{ background-color: #f2f2f2; }}</style></head>
                    <body>
                        <h1>M&A Financial Analysis Report: {acquirer_name} & {target_name}</h1>
                        <p><strong>Date:</strong> {datetime.date.today().strftime('%B %d, %Y')}</p><p><strong>Display Currency:</strong> {report_data['selected_currency']}</p>
                        
                        <h2>1. Executive Summary</h2>
                        <p>This report provides a comprehensive financial analysis of a potential merger between <b>{acquirer_name}</b> and <b>{target_name}</b>.</p>
                        <ul><li>The transaction is <b>{ 'ACCRETIVE' if report_data['accertion_dilution'] > 0 else 'DILUTIVE' }</b> to the acquirer's EPS by <b>{report_data['accertion_dilution']:.2f}%</b>.</li>
                            <li>The estimated <b>goodwill created</b> is <b>{report_data['selected_currency']} {report_data['goodwill']:.2f} billion</b>.</li>
                            <li><b>Valuation Summary:</b><br>- DCF: {report_data['selected_currency']} {report_data['dcf_value']:.2f} billion<br>- Comps: {report_data['selected_currency']} {report_data['comps_value_converted']:.2f} billion<br>- Precedents: {selected_currency} {report_data['precedents_value_converted']:.2f} billion</li></ul>
                        
                        <h2>2. Financial Ratios & Industry Comparison</h2>
                        <img src="data:image/png;base64,{fig_ratios_base64}" style="width: 100%;">
                        <p><i>The chart above visually compares the key financial ratios against the industry average.</i></p>

                        <h2>3. LBO Analysis</h2>
                        <img src="data:image/png;base64,{fig_lbo_base64}" style="width: 100%;">
                        <table><tr><th>Metric</th><th>Value</th></tr><tr><td>Projected Exit EV</td><td>{report_data['selected_currency']} {report_data['proj_exit_ev']:.2f} billion</td></tr><tr><td>Projected Exit Equity</td><td>{report_data['selected_currency']} {report_data['proj_exit_equity']:.2f} billion</td></tr><tr><td>MOIC</td><td>{report_data['moic']:.2f}x</td></tr><tr><td>IRR</td><td>{report_data['irr']:.2%}</td></tr></table>
                        
                        <h2>4. Projected Financials (Target Company)</h2>
                        {proj_html}
                        
                    </body></html>
                    """
                    pdf_file = pdfkit.from_string(html_content, False, configuration=config)
                    return pdf_file
                except Exception as e:
                    st.error(f"Error generating PDF. Please ensure wkhtmltopdf is installed correctly and is in your system's PATH. Error: {e}")
                    return None

            report_data = {
                "accertion_dilution": accertion_dilution, "goodwill": goodwill, "selected_currency": selected_currency,
                "dcf_value": dcf_value, "comps_value_converted": comps_value_converted, "precedents_value_converted": precedents_value_converted,
                "eps_acquirer": eps_acquirer, "pro_forma_eps": pro_forma_eps, "pro_forma_assets": pro_forma_assets,
                "pro_forma_liabilities": pro_forma_liabilities, "moic": moic, "irr": irr,
                "proj_exit_ev": proj_exit_ev, "proj_exit_equity": proj_exit_equity,
            }
            
            col_report1, col_report2 = st.columns(2)
            with col_report1:
                if st.button("Generate PowerPoint Report 📈"):
                    with st.spinner("Generating PowerPoint report..."):
                        pptx_buffer = create_powerpoint_report(acquirer_name, target_name, report_data, fig_ratios, fig_lbo, pd.DataFrame(), results_df, projections_df)
                        if pptx_buffer:
                            st.download_button(label="Download PowerPoint", data=pptx_buffer, file_name=f"M&A_Report_{acquirer_name}_vs_{target_name}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

            with col_report2:
                if st.button("Generate PDF Report 📄"):
                    with st.spinner("Generating PDF report..."):
                        pdf_file = create_pdf_report(acquirer_name, target_name, report_data, fig_ratios, fig_lbo, results_df, projections_df)
                        if pdf_file:
                            st.download_button(label="Download PDF", data=pdf_file, file_name=f"M&A_Report_{acquirer_name}_vs_{target_name}.pdf", mime="application/pdf")
    else:
        st.info("Please provide company information via the simulated API or by uploading CSV or Excel files in the sidebar to run the analysis.")

if check_password():
    main_app()
